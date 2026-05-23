import io
import os
import sys
import uuid
import shutil
import platform
import tempfile
import subprocess
import psycopg2
from typing import List
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import httpx
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
from groq import Groq
from langchain_text_splitters import RecursiveCharacterTextSplitter
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct
from sentence_transformers import SentenceTransformer

from app.core.config import settings
from app.core.storage import get_minio_client

# ── Tesseract path for Windows ────────────────────────────────────────────────
try:
    import pytesseract
    if platform.system() == "Windows":
        pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
except ImportError:
    pass

# ── Resolve yt-dlp, ffmpeg, ffprobe executables ───────────────────────────────
# Priority:
#   1. Same Scripts/ folder as this Python executable (venv install)
#   2. System PATH (shutil.which)
#   3. Hardcoded fallback for known install locations

def _find_exe(name: str) -> str:
    """
    Return the absolute path to an executable, checking:
      1. The venv Scripts/ directory (where pip installs CLI tools)
      2. System PATH
    Raises RuntimeError if not found anywhere.
    """
    # 1. Venv Scripts/ directory (e.g. .venv/Scripts/yt-dlp.exe on Windows)
    venv_scripts = os.path.dirname(sys.executable)  # same dir as python.exe
    candidates = [name, name + ".exe"]
    for cand in candidates:
        venv_path = os.path.join(venv_scripts, cand)
        if os.path.isfile(venv_path):
            return venv_path

    # 2. System PATH
    found = shutil.which(name)
    if found:
        return found

    raise RuntimeError(
        f"'{name}' executable not found.\n"
        f"  - For yt-dlp: run  .venv\\Scripts\\pip install yt-dlp\n"
        f"  - For ffmpeg/ffprobe: winget install --id Gyan.FFmpeg -e  "
        f"then restart your terminal and backend."
    )


# Resolve once at import time so errors surface immediately on startup
try:
    YTDLP_BIN  = _find_exe("yt-dlp")
    FFMPEG_BIN = _find_exe("ffmpeg")
    FFPROBE_BIN = _find_exe("ffprobe")
    print(f"[Ingestion] yt-dlp  : {YTDLP_BIN}")
    print(f"[Ingestion] ffmpeg  : {FFMPEG_BIN}")
    print(f"[Ingestion] ffprobe : {FFPROBE_BIN}")
except RuntimeError as _e:
    # Don't crash the whole backend if ffmpeg isn't installed yet —
    # just warn, and the Whisper path will raise a clear error at call time.
    print(f"[Ingestion] WARNING: {_e}")
    YTDLP_BIN   = "yt-dlp"
    FFMPEG_BIN  = "ffmpeg"
    FFPROBE_BIN = "ffprobe"

# ── Load embedding model ONCE at startup ──────────────────────────────────────
print("[Ingestion] Loading embedding model... (one-time startup cost)")
embed_model = SentenceTransformer("BAAI/bge-base-en-v1.5")
print("[Ingestion] Embedding model loaded.")

EMBED_DIM = 768

qdrant_client = QdrantClient(
    host=settings.QDRANT_HOST,
    port=settings.QDRANT_PORT,
    timeout=120,
    prefer_grpc=False,
)

VECTOR_COLLECTION = "sourcemind_documents"
AUDIO_EXTENSIONS = {"mp3", "wav", "m4a", "ogg", "flac", "webm", "aac"}
IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "tiff", "tif", "bmp", "gif", "webp"}
EXCEL_EXTENSIONS = {"xlsx", "xls", "xlsm", "xlsb"}

MAX_CHUNKS_PER_DOC = 2000
EMBED_BATCH_SIZE = 128
GROQ_WHISPER_MAX_BYTES = 24 * 1024 * 1024   # 24 MB — Groq hard limit is 25 MB


# ─────────────────────────────────────────────────────────────────────────────
#  DB helpers
# ─────────────────────────────────────────────────────────────────────────────

def get_sync_db_conn():
    return psycopg2.connect(
        host=settings.POSTGRES_SERVER,
        port=settings.POSTGRES_PORT,
        user=settings.POSTGRES_USER,
        password=settings.POSTGRES_PASSWORD,
        dbname=settings.POSTGRES_DB
    )


def update_document_status_sync(document_id: str, new_status: str, error_message: str = None):
    try:
        conn = get_sync_db_conn()
        cur = conn.cursor()
        if error_message:
            cur.execute(
                "UPDATE documents SET status = %s, error_message = %s WHERE id = %s",
                (new_status, error_message[:500], document_id)
            )
        else:
            cur.execute(
                "UPDATE documents SET status = %s WHERE id = %s",
                (new_status, document_id)
            )
        conn.commit()
        cur.close()
        conn.close()
    except Exception as e:
        print(f"[DB Update Error] document_id={document_id} status={new_status} error={e}")


# ─────────────────────────────────────────────────────────────────────────────
#  Qdrant
# ─────────────────────────────────────────────────────────────────────────────

def init_qdrant():
    collections = qdrant_client.get_collections().collections
    existing = next((c for c in collections if c.name == VECTOR_COLLECTION), None)
    if existing:
        info = qdrant_client.get_collection(VECTOR_COLLECTION)
        current_size = info.config.params.vectors.size
        if current_size != EMBED_DIM:
            print(f"[Qdrant] Wrong vector size ({current_size} != {EMBED_DIM}). Recreating...")
            qdrant_client.delete_collection(VECTOR_COLLECTION)
            qdrant_client.create_collection(
                collection_name=VECTOR_COLLECTION,
                vectors_config=VectorParams(size=EMBED_DIM, distance=Distance.COSINE),
            )
    else:
        qdrant_client.create_collection(
            collection_name=VECTOR_COLLECTION,
            vectors_config=VectorParams(size=EMBED_DIM, distance=Distance.COSINE),
        )
        print(f"[Qdrant] Collection created with size={EMBED_DIM}")


# ─────────────────────────────────────────────────────────────────────────────
#  File text extractors
# ─────────────────────────────────────────────────────────────────────────────

def extract_text_from_pdf(file_bytes: bytes) -> str:
    text = ""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    for page in doc:
        text += page.get_text() + "\n\n"
    if len(text.strip()) < 100:
        print("[Ingestion] PDF sparse text — attempting OCR fallback")
        text = extract_text_from_pdf_ocr(file_bytes)
    return text


def extract_text_from_pdf_ocr(file_bytes: bytes) -> str:
    try:
        import pytesseract
        from pdf2image import convert_from_bytes
        poppler_path = None
        if platform.system() == "Windows":
            poppler_path = r"C:\poppler\poppler-24.02.0\Library\bin"
            pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        pages = convert_from_bytes(file_bytes, dpi=200, poppler_path=poppler_path)
        text_parts = []
        for i, page_img in enumerate(pages):
            text_parts.append(f"--- Page {i+1} ---\n{pytesseract.image_to_string(page_img, lang='eng')}")
            print(f"[OCR] Page {i+1}/{len(pages)} done")
        return "\n\n".join(text_parts)
    except ImportError:
        raise ValueError("OCR requires: pip install pytesseract pdf2image Pillow")
    except Exception as e:
        raise ValueError(f"OCR failed: {e}")


def extract_text_from_image(file_bytes: bytes) -> str:
    try:
        import pytesseract
        from PIL import Image
        if platform.system() == "Windows":
            pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        return pytesseract.image_to_string(Image.open(io.BytesIO(file_bytes)), lang="eng")
    except ImportError:
        raise ValueError("Image OCR requires: pip install pytesseract Pillow")
    except Exception as e:
        raise ValueError(f"Image OCR failed: {e}")


def extract_text_from_excel(file_bytes: bytes, file_ext: str) -> str:
    import openpyxl
    workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    text_parts = []
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        text_parts.append(f"=== Sheet: {sheet_name} ===")
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            text_parts.append("(empty sheet)")
            continue
        headers = [str(c) if c is not None else "" for c in rows[0]]
        if any(h.strip() for h in headers):
            text_parts.append(" | ".join(headers))
            text_parts.append("-" * 60)
            data_rows = rows[1:]
        else:
            data_rows = rows
        for row in data_rows:
            if all(cell is None or str(cell).strip() == "" for cell in row):
                continue
            text_parts.append(" | ".join(str(c) if c is not None else "" for c in row))
        text_parts.append("")
    workbook.close()
    return "\n".join(text_parts)


def extract_text_from_csv(file_bytes: bytes) -> str:
    import csv
    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            content = file_bytes.decode(encoding)
            rows = list(csv.reader(content.splitlines()))
            return "\n".join(" | ".join(row) for row in rows)
        except (UnicodeDecodeError, csv.Error):
            continue
    return file_bytes.decode("utf-8", errors="ignore")


def extract_text_from_docx(file_bytes: bytes) -> str:
    doc = docx.Document(io.BytesIO(file_bytes))
    parts = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(parts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for i, slide in enumerate(prs.slides):
        parts.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts)


def extract_text_from_url(url: str) -> str:
    headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
    with httpx.Client(headers=headers, timeout=15.0, follow_redirects=True) as client:
        response = client.get(url)
        response.raise_for_status()
    soup = BeautifulSoup(response.text, "html.parser")
    for el in soup(["script", "style", "head", "nav", "footer", "iframe", "aside"]):
        el.decompose()
    lines = (line.strip() for line in soup.get_text(separator="\n").splitlines())
    return "\n".join(line for line in lines if line)


# ─────────────────────────────────────────────────────────────────────────────
#  YouTube: transcript API → Groq Whisper fallback
# ─────────────────────────────────────────────────────────────────────────────

class _TranscriptUnavailable(Exception):
    """Signals that the transcript API path failed — triggers Whisper fallback."""
    pass


def _extract_video_id(url: str) -> str:
    if "v=" in url:
        return url.split("v=")[1].split("&")[0]
    if "youtu.be/" in url:
        return url.split("youtu.be/")[-1].split("?")[0]
    raise ValueError(f"Could not extract video ID from URL: {url}")


def _transcribe_audio_bytes_groq(audio_bytes: bytes, filename: str = "audio.mp3") -> str:
    """Send raw audio bytes to Groq whisper-large-v3. File must be < 25 MB."""
    if not settings.GROQ_API_KEY:
        raise ValueError("GROQ_API_KEY is not configured")
    client = Groq(api_key=settings.GROQ_API_KEY)
    buf = io.BytesIO(audio_bytes)
    buf.name = filename
    result = client.audio.transcriptions.create(
        file=buf,
        model="whisper-large-v3",
        response_format="text",
    )
    return result if isinstance(result, str) else result.text


def _split_audio_file(audio_path: str, max_bytes: int = GROQ_WHISPER_MAX_BYTES) -> list:
    """
    Split audio into time-based chunks, each guaranteed under max_bytes.

    Strategy:
      - Use WAV (PCM) output so every chunk has a self-contained, valid
        header with NO codec dependencies (no libmp3lame needed).
      - WAV files are larger per second, so we calculate chunk_duration
        based on the ACTUAL WAV bitrate of the source, not the compressed
        MP3 size, to stay safely under the Groq 25 MB limit.
      - Groq accepts WAV natively.

    Returns list of .wav file paths in a fresh temp dir.
    The caller (extract_text_from_youtube_whisper) already wraps everything
    in a try/finally that calls shutil.rmtree on the PARENT tmp dir, so
    these chunk files are cleaned up automatically.
    """
    file_size = os.path.getsize(audio_path)

    # ── Step 1: probe duration ────────────────────────────────────────────
    probe = subprocess.run(
        [
            FFPROBE_BIN, "-v", "error",
            "-show_entries", "format=duration",
            "-of", "default=noprint_wrappers=1:nokey=1",
            audio_path,
        ],
        capture_output=True, text=True, check=True,
    )
    total_duration = float(probe.stdout.strip())

    # ── Step 2: estimate WAV output size ─────────────────────────────────
    # PCM 16-bit mono 16 kHz (Groq downsamples anyway) = 32 000 bytes/sec
    # Use 16 kHz mono so chunks are small and always valid.
    WAV_BYTES_PER_SEC = 16000 * 2  # 16 kHz, 16-bit mono = 32 000 B/s
    safe_limit = int(max_bytes * 0.90)           # 90 % of 24 MB = 21.6 MB
    chunk_duration = safe_limit / WAV_BYTES_PER_SEC  # seconds per chunk

    n_chunks = int(total_duration / chunk_duration) + 1

    if n_chunks == 1:
        # File fits in one chunk — just convert the whole thing to WAV
        tmp_dir = tempfile.mkdtemp()
        out_path = os.path.join(tmp_dir, "chunk_000.wav")
        subprocess.run(
            [
                FFMPEG_BIN, "-y", "-i", audio_path,
                "-ar", "16000", "-ac", "1", "-sample_fmt", "s16",
                out_path,
            ],
            capture_output=True, check=True,
        )
        print(f"[Whisper] Single WAV chunk: {os.path.getsize(out_path)/1024/1024:.1f} MB")
        return [out_path]

    # ── Step 3: split into N WAV chunks ──────────────────────────────────
    tmp_dir = tempfile.mkdtemp()
    chunk_paths = []

    for i in range(n_chunks):
        start = i * chunk_duration
        out_path = os.path.join(tmp_dir, f"chunk_{i:03d}.wav")
        result = subprocess.run(
            [
                FFMPEG_BIN, "-y",
                "-ss", str(start),
                "-t",  str(chunk_duration),
                "-i",  audio_path,
                "-ar", "16000",        # 16 kHz — Groq's native rate
                "-ac", "1",            # mono
                "-sample_fmt", "s16",  # 16-bit PCM
                out_path,
            ],
            capture_output=True, text=True,
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"ffmpeg chunk {i} failed (exit {result.returncode}).\n"
                f"stderr: {result.stderr[-800:]}"
            )
        chunk_size = os.path.getsize(out_path) if os.path.exists(out_path) else 0
        if chunk_size < 1000:
            # Empty or near-empty tail chunk — skip it
            os.remove(out_path)
            continue
        chunk_paths.append(out_path)
        print(f"[Whisper] Chunk {i+1}: {chunk_size/1024/1024:.1f} MB ({start:.0f}s–{start+chunk_duration:.0f}s)")

    if not chunk_paths:
        raise RuntimeError("ffmpeg produced zero valid WAV chunks.")

    print(f"[Whisper] Split into {len(chunk_paths)} WAV chunk(s) "
          f"(source {file_size/1024/1024:.1f} MB)")
    return chunk_paths


def extract_text_from_youtube_whisper(url: str) -> str:
    """
    Download YouTube audio with yt-dlp, split into <24 MB chunks,
    transcribe each with Groq whisper-large-v3, then join.
    """
    if not settings.GROQ_API_KEY:
        raise ValueError("GROQ_API_KEY is not configured — cannot use Whisper fallback")

    video_id = _extract_video_id(url)
    tmp_dir = tempfile.mkdtemp()
    # Use %(ext)s so yt-dlp fills in the actual extension it produces
    output_template = os.path.join(tmp_dir, f"{video_id}.%(ext)s")

    print(f"[Whisper] Downloading audio for video_id={video_id} ...")
    print(f"[Whisper] Using yt-dlp: {YTDLP_BIN}")
    print(f"[Whisper] Using ffmpeg: {FFMPEG_BIN}")

    try:
        result = subprocess.run(
            [
                YTDLP_BIN,
                "--no-playlist",
                "-x",
                "--audio-format", "mp3",
                "--audio-quality", "5",      # ~128 kbps — fine for speech
                "--ffmpeg-location", os.path.dirname(FFMPEG_BIN),  # tell yt-dlp where ffmpeg is
                "-o", output_template,
                "--no-warnings",
                url,
            ],
            capture_output=True, text=True,
        )

        if result.returncode != 0:
            error_detail = (result.stderr or result.stdout or "").strip()
            raise ValueError(
                f"yt-dlp failed for video ({video_id}). Detail: {error_detail}"
            )

        # Find the downloaded file (yt-dlp fills in the extension)
        audio_path = None
        for f in os.listdir(tmp_dir):
            if f.startswith(video_id):
                audio_path = os.path.join(tmp_dir, f)
                break

        if not audio_path or not os.path.exists(audio_path):
            raise ValueError(
                f"yt-dlp completed but output file not found for video ({video_id}). "
                "Ensure ffmpeg is installed and on PATH."
            )

        file_size_mb = os.path.getsize(audio_path) / 1024 / 1024
        print(f"[Whisper] Downloaded {file_size_mb:.1f} MB → {audio_path}")

        # Split if needed
        chunk_paths = _split_audio_file(audio_path)

        # Transcribe each chunk
        transcript_parts = []
        for i, chunk_path in enumerate(chunk_paths):
            print(f"[Whisper] Transcribing chunk {i+1}/{len(chunk_paths)} ...")
            with open(chunk_path, "rb") as f:
                audio_bytes = f.read()
            chunk_text = _transcribe_audio_bytes_groq(
                audio_bytes,
                filename=os.path.basename(chunk_path),
            )
            transcript_parts.append(chunk_text.strip())
            print(f"[Whisper] Chunk {i+1} → {len(chunk_text):,} chars")

        full_transcript = "\n\n".join(p for p in transcript_parts if p)

        if not full_transcript.strip():
            raise ValueError(
                f"Groq Whisper returned an empty transcript for video ({video_id}). "
                "The audio may contain no speech."
            )

        print(f"[Whisper] Done: {len(full_transcript):,} chars for video_id={video_id}")
        return full_transcript

    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def extract_text_from_youtube(url: str) -> str:
    """
    Primary  : youtube_transcript_api (instant, no download).
    Fallback : yt-dlp + Groq whisper-large-v3.
    """
    import time

    video_id = _extract_video_id(url)
    print(f"[YouTube] Trying transcript API for video_id={video_id}")

    try:
        transcript_list = None
        for attempt in range(3):
            try:
                transcript_list = YouTubeTranscriptApi.list_transcripts(video_id)
                break
            except Exception as e:
                err_str = str(e)
                if "429" in err_str or "Too Many Requests" in err_str:
                    wait = (attempt + 1) * 15
                    print(f"[YouTube] Rate limited, waiting {wait}s (attempt {attempt+1}/3)...")
                    time.sleep(wait)
                elif "403" in err_str or "Forbidden" in err_str:
                    raise _TranscriptUnavailable(f"403 from YouTube for video ({video_id})")
                else:
                    raise _TranscriptUnavailable(str(e))

        if transcript_list is None:
            raise _TranscriptUnavailable(f"Rate-limited after 3 attempts for video ({video_id})")

        all_transcripts = []
        try:
            for t in transcript_list:
                all_transcripts.append(t)
        except Exception:
            pass

        if not all_transcripts:
            raise _TranscriptUnavailable(f"No transcripts via API for video ({video_id})")

        # Fetch order: English → manual → auto-generated
        fetch_attempts = []
        for t in all_transcripts:
            if t.language_code in ("en", "en-US", "en-GB"):
                fetch_attempts.append((t, t.language))
        for t in all_transcripts:
            if not t.is_generated and t.language_code not in ("en", "en-US", "en-GB"):
                fetch_attempts.append((t, f"{t.language} (manual)"))
        for t in all_transcripts:
            if t.is_generated and t.language_code not in ("en", "en-US", "en-GB"):
                fetch_attempts.append((t, f"{t.language} (auto)"))
        if not fetch_attempts:
            fetch_attempts = [(all_transcripts[0], all_transcripts[0].language)]

        for transcript_obj, lang_label in fetch_attempts:
            for attempt in range(2):
                try:
                    entries = transcript_obj.fetch()
                    texts = []
                    for entry in entries:
                        texts.append(entry.get("text", "") if isinstance(entry, dict)
                                     else getattr(entry, "text", str(entry)))
                    full_text = "\n".join(t for t in texts if t.strip())
                    if full_text.strip():
                        print(f"[YouTube] Transcript API success: {len(full_text):,} chars ({lang_label})")
                        return full_text
                    break
                except Exception as e:
                    err_str = str(e)
                    if "429" in err_str or "Too Many Requests" in err_str:
                        time.sleep((attempt + 1) * 15)
                    elif "no element found" in err_str or "line 1, column 0" in err_str:
                        break
                    elif "403" in err_str:
                        raise _TranscriptUnavailable(f"403 on fetch for video ({video_id})")
                    else:
                        break

        raise _TranscriptUnavailable(f"All transcript API results empty for video ({video_id})")

    except _TranscriptUnavailable as e:
        print(f"[YouTube] Transcript API unavailable: {e}")
        print(f"[YouTube] Falling back to yt-dlp + Groq Whisper for video_id={video_id}")
        return extract_text_from_youtube_whisper(url)


def extract_text_from_audio(file_bytes: bytes, file_ext: str) -> str:
    """Transcribe an uploaded audio file via Groq whisper-large-v3."""
    if not settings.GROQ_API_KEY:
        raise ValueError("GROQ_API_KEY is not configured")
    client = Groq(api_key=settings.GROQ_API_KEY)
    filename = f"audio.{file_ext}" if file_ext else "audio.mp3"
    audio_file = io.BytesIO(file_bytes)
    audio_file.name = filename
    return client.audio.transcriptions.create(
        file=audio_file,
        model="whisper-large-v3",
        response_format="text"
    )


# ─────────────────────────────────────────────────────────────────────────────
#  Main ingestion pipeline
# ─────────────────────────────────────────────────────────────────────────────

def process_document(document_id: str, file_path: str, file_type: str, notebook_id: str):
    """Parse → chunk → embed → store in Qdrant."""
    import numpy as np
    print(f"[Ingestion] Starting: document_id={document_id} file_type={file_type}")
    update_document_status_sync(document_id, "processing")

    try:
        file_type_lower = file_type.lower()
        text = ""

        # 1. Extract text
        if file_type_lower in ["url", "youtube"]:
            if file_type_lower == "url":
                text = extract_text_from_url(file_path)
            else:
                text = extract_text_from_youtube(file_path)
        else:
            minio_client = get_minio_client()
            bucket_name = "sourcemind-documents"
            object_name = file_path.split(f"{bucket_name}/")[1]
            response = minio_client.get_object(bucket_name, object_name)
            try:
                file_bytes = response.read()
            finally:
                response.close()
                response.release_conn()

            if file_type_lower == "pdf":
                text = extract_text_from_pdf(file_bytes)
            elif file_type_lower in ["docx", "doc"]:
                text = extract_text_from_docx(file_bytes)
            elif file_type_lower in ["pptx", "ppt"]:
                text = extract_text_from_pptx(file_bytes)
            elif file_type_lower in AUDIO_EXTENSIONS:
                text = extract_text_from_audio(file_bytes, file_type_lower)
            elif file_type_lower in IMAGE_EXTENSIONS:
                text = extract_text_from_image(file_bytes)
            elif file_type_lower in EXCEL_EXTENSIONS:
                text = extract_text_from_excel(file_bytes, file_type_lower)
            elif file_type_lower == "csv":
                text = extract_text_from_csv(file_bytes)
            else:
                text = file_bytes.decode("utf-8", errors="ignore")

        if not text.strip():
            raise ValueError("Extracted text is empty")

        print(f"[Ingestion] Extracted {len(text):,} chars from document_id={document_id}")

        # 2. Chunk
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1500,
            chunk_overlap=300,
            length_function=len,
            separators=["\n\n\n", "\n\n", "\n", ". ", " ", ""]
        )
        chunks = text_splitter.split_text(text)
        if not chunks:
            raise ValueError("No chunks generated from text")

        total_chunks = len(chunks)
        print(f"[Ingestion] Created {total_chunks:,} chunks for document_id={document_id}")

        if total_chunks > MAX_CHUNKS_PER_DOC:
            step = total_chunks / MAX_CHUNKS_PER_DOC
            indices = [int(i * step) for i in range(MAX_CHUNKS_PER_DOC)]
            chunks = [chunks[i] for i in indices]
            print(f"[Ingestion] Sampled {MAX_CHUNKS_PER_DOC:,} from {total_chunks:,} chunks")

        # 3. Embed
        all_embeddings = []
        n_chunks = len(chunks)
        for batch_start in range(0, n_chunks, EMBED_BATCH_SIZE):
            batch = chunks[batch_start:batch_start + EMBED_BATCH_SIZE]
            batch_emb = embed_model.encode(
                batch, convert_to_numpy=True,
                show_progress_bar=False, batch_size=EMBED_BATCH_SIZE
            )
            all_embeddings.append(batch_emb)
            done = min(batch_start + EMBED_BATCH_SIZE, n_chunks)
            print(f"[Ingestion] Embedded {done}/{n_chunks} ({100*done//n_chunks}%) document_id={document_id}")

        embeddings = np.vstack(all_embeddings)

        # 4. Store in Qdrant
        points = [
            PointStruct(
                id=str(uuid.uuid4()),
                vector=embedding.tolist(),
                payload={
                    "document_id": str(document_id),
                    "notebook_id": str(notebook_id),
                    "text": chunk_text,
                    "chunk_index": i,
                }
            )
            for i, (chunk_text, embedding) in enumerate(zip(chunks, embeddings))
        ]

        init_qdrant()
        fresh_qdrant = QdrantClient(
            host=settings.QDRANT_HOST,
            port=settings.QDRANT_PORT,
            timeout=120,
            prefer_grpc=False,
        )
        UPSERT_BATCH = 200
        for i in range(0, len(points), UPSERT_BATCH):
            batch = points[i:i + UPSERT_BATCH]
            for attempt in range(3):
                try:
                    fresh_qdrant.upsert(collection_name=VECTOR_COLLECTION, points=batch)
                    break
                except Exception as upsert_err:
                    if attempt == 2:
                        raise
                    import time
                    print(f"[Ingestion] Upsert attempt {attempt+1} failed: {upsert_err}, retrying...")
                    time.sleep(3)
            print(f"[Ingestion] Upserted {min(i+UPSERT_BATCH, len(points))}/{len(points)} vectors")

        # 5. Done
        update_document_status_sync(document_id, "done")
        print(f"[Ingestion] DONE: document_id={document_id}")

    except Exception as e:
        print(f"[Ingestion] ERROR: document_id={document_id} error={e}")
        update_document_status_sync(document_id, "error", error_message=str(e))
